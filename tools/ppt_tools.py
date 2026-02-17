"""
=============================================================================
 ULTRA PPT ENGINE v2.0 — Professional Presentation Generator
 Supports: gradient backgrounds, decorative shapes, native charts, card layouts,
 preset themes, timeline, progress bars, stat numbers, icons, and more.
=============================================================================
"""
import os
import copy
import math
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn, nsmap
from .base import get_workspace_path

# ============================================================================
#  PRESET THEMES — Professional color palettes
# ============================================================================

THEMES = {
    "ocean": {
        "primary": "#1A5276", "secondary": "#2980B9", "accent": "#3498DB",
        "light": "#D6EAF8", "dark": "#0E2F44", "text_dark": "#1C2833",
        "text_light": "#FFFFFF", "bg": "#F0F8FF", "gradient_start": "#1A5276",
        "gradient_end": "#2980B9", "card_bg": "#FFFFFF", "card_border": "#AED6F1",
    },
    "forest": {
        "primary": "#1E8449", "secondary": "#27AE60", "accent": "#2ECC71",
        "light": "#D5F5E3", "dark": "#0B3D1E", "text_dark": "#1C2833",
        "text_light": "#FFFFFF", "bg": "#F0FFF0", "gradient_start": "#1E8449",
        "gradient_end": "#27AE60", "card_bg": "#FFFFFF", "card_border": "#ABEBC6",
    },
    "sunset": {
        "primary": "#C0392B", "secondary": "#E74C3C", "accent": "#F39C12",
        "light": "#FDEDEC", "dark": "#641E16", "text_dark": "#1C2833",
        "text_light": "#FFFFFF", "bg": "#FFF5F0", "gradient_start": "#C0392B",
        "gradient_end": "#E74C3C", "card_bg": "#FFFFFF", "card_border": "#F5B7B1",
    },
    "royal": {
        "primary": "#6C3483", "secondary": "#8E44AD", "accent": "#BB8FCE",
        "light": "#F4ECF7", "dark": "#3B1A5C", "text_dark": "#1C2833",
        "text_light": "#FFFFFF", "bg": "#FAF0FF", "gradient_start": "#6C3483",
        "gradient_end": "#8E44AD", "card_bg": "#FFFFFF", "card_border": "#D2B4DE",
    },
    "midnight": {
        "primary": "#2C3E50", "secondary": "#34495E", "accent": "#1ABC9C",
        "light": "#EBF5FB", "dark": "#1B2631", "text_dark": "#1C2833",
        "text_light": "#FFFFFF", "bg": "#F8F9FA", "gradient_start": "#2C3E50",
        "gradient_end": "#34495E", "card_bg": "#FFFFFF", "card_border": "#AEB6BF",
    },
    "coral": {
        "primary": "#E8725C", "secondary": "#F09E8C", "accent": "#F7C59F",
        "light": "#FFF0EB", "dark": "#8B3A2F", "text_dark": "#2D2D2D",
        "text_light": "#FFFFFF", "bg": "#FFFAF8", "gradient_start": "#E8725C",
        "gradient_end": "#F09E8C", "card_bg": "#FFFFFF", "card_border": "#F5C6BA",
    },
    "tech": {
        "primary": "#0D47A1", "secondary": "#1565C0", "accent": "#42A5F5",
        "light": "#E3F2FD", "dark": "#0A1929", "text_dark": "#1A237E",
        "text_light": "#FFFFFF", "bg": "#F5F9FF", "gradient_start": "#0D47A1",
        "gradient_end": "#1565C0", "card_bg": "#FFFFFF", "card_border": "#90CAF9",
    },
    "elegant": {
        "primary": "#2C2C2C", "secondary": "#555555", "accent": "#C9A96E",
        "light": "#F5F0E8", "dark": "#1A1A1A", "text_dark": "#2C2C2C",
        "text_light": "#FFFFFF", "bg": "#FAFAF5", "gradient_start": "#2C2C2C",
        "gradient_end": "#555555", "card_bg": "#FFFFFF", "card_border": "#D5C4A1",
    },
}

# ============================================================================
#  UTILITY FUNCTIONS
# ============================================================================

def _hex(color):
    """Convert hex color string to RGBColor."""
    c = color.lstrip('#')
    if len(c) == 3:
        c = ''.join([ch * 2 for ch in c])
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def _inches(val, total=None):
    """Parse dimension value — supports float (inches), string percent, or Emu."""
    if isinstance(val, str) and val.endswith('%'):
        return Inches(float(val.rstrip('%')) / 100 * (total or 13.333))
    return Inches(float(val))


def _pt(val):
    return Pt(val)


def _emu(inches_val):
    return Emu(int(inches_val * 914400))


def _darken(hex_color, factor=0.7):
    """Darken a hex color by a factor."""
    c = hex_color.lstrip('#')
    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
    r, g, b = int(r * factor), int(g * factor), int(b * factor)
    return f"#{r:02x}{g:02x}{b:02x}"


def _lighten(hex_color, factor=0.3):
    """Lighten a hex color by mixing with white."""
    c = hex_color.lstrip('#')
    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
    r = int(r + (255 - r) * factor)
    g = int(g + (255 - g) * factor)
    b = int(b + (255 - b) * factor)
    return f"#{r:02x}{g:02x}{b:02x}"


def _get_theme(name="midnight"):
    """Get theme colors by name, with fallback to midnight."""
    return THEMES.get(name, THEMES["midnight"])


# ============================================================================
#  GRADIENT FILL via XML (supports multi-stop linear gradients)
# ============================================================================

def _apply_gradient_fill(fill_target, color_start, color_end, angle=270, extra_stops=None):
    """
    Apply a linear gradient fill to a shape or slide background via direct XML manipulation.
    fill_target: the shape or slide background element that has .fill
    angle: degrees (0=left-to-right, 90=bottom-to-top, 270=top-to-bottom)
    extra_stops: list of {"pos": 0-100, "color": "#hex"} for multi-stop gradients
    """
    fill = fill_target.fill
    fill.gradient()
    fill.gradient_angle = angle

    stops = fill.gradient_stops
    # Set first stop
    stops[0].color.rgb = _hex(color_start)
    stops[0].position = 0.0
    # Set last stop
    stops[1].color.rgb = _hex(color_end)
    stops[1].position = 1.0

    # Add extra stops via XML if needed
    if extra_stops:
        gs_lst = fill._fill.findall(qn('a:gsLst'))
        if gs_lst:
            gs_lst = gs_lst[0]
        else:
            # Try direct child
            for child in fill._fill:
                if child.tag.endswith('gsLst'):
                    gs_lst = child
                    break
        if gs_lst is not None:
            for es in extra_stops:
                pos = int(es["pos"] * 1000)  # Convert percentage to 1/1000ths
                gs_el = etree.SubElement(gs_lst, qn('a:gs'))
                gs_el.set('pos', str(pos))
                srgb = etree.SubElement(gs_el, qn('a:srgbClr'))
                srgb.set('val', es["color"].lstrip('#'))


def _apply_solid_fill(shape, color):
    """Apply solid color fill to a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)


def _apply_bg_gradient(slide, color_start, color_end, angle=270):
    """Apply gradient background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].color.rgb = _hex(color_start)
    stops[0].position = 0.0
    stops[1].color.rgb = _hex(color_end)
    stops[1].position = 1.0


def _apply_bg_solid(slide, color):
    """Apply solid background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex(color)


# ============================================================================
#  TEXT HELPERS — Rich text with multiple runs, paragraphs, bullet points
# ============================================================================

def _set_run_style(run, font_name="微软雅黑", font_size=18, color="#333333",
                   bold=False, italic=False, underline=False):
    """Style a single text run."""
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = _hex(color)
    run.font.bold = bold
    run.font.italic = italic
    if underline:
        run.font.underline = True


def _add_text_to_frame(tf, text_config, default_font="微软雅黑"):
    """
    Add rich text to a text frame. Supports:
    - Simple string: {"text": "Hello"}
    - Paragraphs list: {"paragraphs": [{"text": "line1", "bold": true}, ...]}
    - Bullet list: {"bullets": ["item1", "item2"]}
    """
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    margin_l = text_config.get("margin_left", 10)
    margin_r = text_config.get("margin_right", 10)
    margin_t = text_config.get("margin_top", 5)
    margin_b = text_config.get("margin_bottom", 5)
    tf.margin_left = Pt(margin_l)
    tf.margin_right = Pt(margin_r)
    tf.margin_top = Pt(margin_t)
    tf.margin_bottom = Pt(margin_b)

    v_anchor = text_config.get("vertical_anchor", "TOP")
    tf.vertical_anchor = getattr(MSO_ANCHOR, v_anchor.upper(), MSO_ANCHOR.TOP)

    font_size = text_config.get("font_size", 18)
    color = text_config.get("color", "#333333")
    bold = text_config.get("bold", False)
    italic = text_config.get("italic", False)
    align_str = text_config.get("align", "LEFT")
    alignment = getattr(PP_ALIGN, align_str.upper(), PP_ALIGN.LEFT)
    line_spacing = text_config.get("line_spacing", 1.2)

    # Simple text
    if "text" in text_config and not isinstance(text_config["text"], list):
        p = tf.paragraphs[0]
        p.alignment = alignment
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = str(text_config["text"])
        _set_run_style(run, default_font, font_size, color, bold, italic,
                       text_config.get("underline", False))
        return

    # Paragraphs list (rich multi-paragraph text)
    if "paragraphs" in text_config:
        for i, para_cfg in enumerate(text_config["paragraphs"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = getattr(PP_ALIGN, para_cfg.get("align", align_str).upper(), alignment)
            p.line_spacing = para_cfg.get("line_spacing", line_spacing)
            p.space_before = Pt(para_cfg.get("space_before", 4))
            p.space_after = Pt(para_cfg.get("space_after", 4))

            # Support runs within a paragraph
            if "runs" in para_cfg:
                for run_cfg in para_cfg["runs"]:
                    run = p.add_run()
                    run.text = str(run_cfg.get("text", ""))
                    _set_run_style(
                        run,
                        run_cfg.get("font_name", default_font),
                        run_cfg.get("font_size", font_size),
                        run_cfg.get("color", color),
                        run_cfg.get("bold", bold),
                        run_cfg.get("italic", italic),
                        run_cfg.get("underline", False),
                    )
            else:
                run = p.add_run()
                run.text = str(para_cfg.get("text", ""))
                _set_run_style(
                    run,
                    para_cfg.get("font_name", default_font),
                    para_cfg.get("font_size", font_size),
                    para_cfg.get("color", color),
                    para_cfg.get("bold", bold),
                    para_cfg.get("italic", italic),
                    para_cfg.get("underline", False),
                )
        return

    # Bullet list
    if "bullets" in text_config:
        for i, item in enumerate(text_config["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = alignment
            p.line_spacing = line_spacing
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            p.level = 0

            bullet_text = item if isinstance(item, str) else item.get("text", "")
            bullet_level = 0 if isinstance(item, str) else item.get("level", 0)
            p.level = bullet_level

            # Add bullet character
            run = p.add_run()
            prefix = "  " * bullet_level + "•  "
            run.text = prefix + bullet_text
            item_color = color if isinstance(item, str) else item.get("color", color)
            item_bold = bold if isinstance(item, str) else item.get("bold", bold)
            item_size = font_size if isinstance(item, str) else item.get("font_size", font_size)
            _set_run_style(run, default_font, item_size, item_color, item_bold, italic)
        return


# ============================================================================
#  SHAPE BUILDERS — Decorative elements, cards, icons
# ============================================================================

def _add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None,
                      border_width=0, corner_radius=None, shadow=False):
    """Add a rounded rectangle shape with optional styling."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _inches(left), _inches(top),
        _inches(width), _inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(fill_color)

    if border_color and border_width > 0:
        shape.line.color.rgb = _hex(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    if shadow:
        shape.shadow.inherit = False

    return shape


def _add_circle(slide, left, top, size, fill_color, border_color=None, border_width=0):
    """Add a circle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, _inches(left), _inches(top),
        _inches(size), _inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(fill_color)
    if border_color and border_width > 0:
        shape.line.color.rgb = _hex(border_color)
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def _add_line(slide, start_x, start_y, end_x, end_y, color, width=2):
    """Add a line shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        _inches(start_x), _inches(start_y),
        _inches(end_x - start_x) if end_x > start_x else Inches(0.01),
        _inches(end_y - start_y) if end_y > start_y else Inches(0.01),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)
    shape.line.fill.background()
    return shape


def _add_decorative_bar(slide, x, y, w, h, color):
    """Add a thin decorative bar/line."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _inches(x), _inches(y), _inches(w), _inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex(color)
    bar.line.fill.background()
    return bar


def _add_text_box(slide, left, top, width, height, text_config, default_font="微软雅黑"):
    """Add a text box with rich text support."""
    txBox = slide.shapes.add_textbox(_inches(left), _inches(top), _inches(width), _inches(height))
    _add_text_to_frame(txBox.text_frame, text_config, default_font)
    return txBox


# ============================================================================
#  NATIVE CHART SUPPORT (embedded PowerPoint charts)
# ============================================================================

def _add_native_chart(slide, chart_config, left=1, top=2, width=8, height=4.5):
    """
    Add a native PowerPoint chart to a slide.
    chart_config:
      - chart_type: "column", "bar", "line", "pie", "area", "scatter", "doughnut"
      - categories: ["Cat1", "Cat2", ...]
      - series: [{"name": "S1", "values": [10, 20, ...]}, ...]
      - title: chart title (optional)
      - show_legend: bool
      - show_data_labels: bool
      - colors: ["#hex1", "#hex2", ...] for series colors
      - legend_position: "BOTTOM", "RIGHT", etc.
    """
    chart_type_str = chart_config.get("chart_type", "column").upper()
    type_map = {
        "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "COLUMN_STACKED": XL_CHART_TYPE.COLUMN_STACKED,
        "BAR": XL_CHART_TYPE.BAR_CLUSTERED,
        "BAR_STACKED": XL_CHART_TYPE.BAR_STACKED,
        "LINE": XL_CHART_TYPE.LINE_MARKERS,
        "LINE_SMOOTH": XL_CHART_TYPE.LINE_MARKERS_STACKED,
        "PIE": XL_CHART_TYPE.PIE,
        "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
        "AREA": XL_CHART_TYPE.AREA,
        "AREA_STACKED": XL_CHART_TYPE.AREA_STACKED,
        "SCATTER": XL_CHART_TYPE.XY_SCATTER,
    }
    xl_type = type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    categories = chart_config.get("categories", [])
    series_data = chart_config.get("series", [])

    if chart_type_str == "SCATTER":
        chart_data = XyChartData()
        for s in series_data:
            series = chart_data.add_series(s.get("name", "Series"))
            x_vals = s.get("x_values", [])
            y_vals = s.get("values", s.get("y_values", []))
            for xv, yv in zip(x_vals, y_vals):
                series.add_data_point(xv, yv)
    else:
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in series_data:
            chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    x, y = _inches(left), _inches(top)
    cx, cy = _inches(width), _inches(height)
    graphic_frame = slide.shapes.add_chart(xl_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart

    # Title
    if chart_config.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_config["title"]
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    else:
        chart.has_title = False

    # Legend
    if chart_config.get("show_legend", True) and len(series_data) > 1:
        chart.has_legend = True
        legend_pos = chart_config.get("legend_position", "BOTTOM").upper()
        chart.legend.position = getattr(XL_LEGEND_POSITION, legend_pos, XL_LEGEND_POSITION.BOTTOM)
        chart.legend.include_in_layout = False
    elif not chart_config.get("show_legend", True):
        chart.has_legend = False

    # Data labels
    if chart_config.get("show_data_labels", False):
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.number_format = chart_config.get("number_format", "0")

    # Series colors
    colors = chart_config.get("colors", [])
    if colors:
        for i, series in enumerate(chart.series):
            if i < len(colors):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = _hex(colors[i])

    return graphic_frame


# ============================================================================
#  TABLE BUILDER — Professional styled tables
# ============================================================================

def _add_styled_table(slide, table_config, theme, default_font="微软雅黑"):
    """
    Add a professionally styled table.
    table_config:
      - left, top, width, height: positioning
      - data: 2D array (first row = headers)
      - header_bg: header background color
      - header_fg: header text color
      - stripe_colors: [color1, color2] for alternating rows
      - font_size: text size
      - header_font_size: header text size
    """
    data = table_config.get("data", [])
    if not data:
        return None

    rows = len(data)
    cols = len(data[0]) if data else 0
    if rows == 0 or cols == 0:
        return None

    left = _inches(table_config.get("left", 1))
    top = _inches(table_config.get("top", 2))
    width = _inches(table_config.get("width", 11))
    height = _inches(table_config.get("height", rows * 0.5))

    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    header_bg = table_config.get("header_bg", theme["primary"])
    header_fg = table_config.get("header_fg", "#FFFFFF")
    stripe1 = "#FFFFFF"
    stripe2 = theme.get("light", "#F5F5F5")
    stripes = table_config.get("stripe_colors", [stripe1, stripe2])
    font_size = table_config.get("font_size", 12)
    header_font_size = table_config.get("header_font_size", 13)

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            if c_idx >= cols:
                continue
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)

            # Style
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            if r_idx == 0:
                # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex(header_bg)
                for run in p.runs:
                    run.font.size = Pt(header_font_size)
                    run.font.bold = True
                    run.font.color.rgb = _hex(header_fg)
                    run.font.name = default_font
                if not p.runs:
                    p.font.size = Pt(header_font_size)
                    p.font.bold = True
                    p.font.color.rgb = _hex(header_fg)
                    p.font.name = default_font
            else:
                # Data rows with striping
                stripe_color = stripes[r_idx % len(stripes)] if stripes else "#FFFFFF"
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex(stripe_color)
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = _hex(theme.get("text_dark", "#333333"))
                    run.font.name = default_font
                if not p.runs:
                    p.font.size = Pt(font_size)
                    p.font.color.rgb = _hex(theme.get("text_dark", "#333333"))
                    p.font.name = default_font

    return shape


# ============================================================================
#  HIGH-LEVEL SLIDE BUILDERS — Each creates a complete, beautifully designed slide
# ============================================================================

def _build_title_slide(prs, slide_data, theme, default_font):
    """
    Build a professional title/cover slide with gradient background and decorative elements.
    slide_data:
      - title: main title text
      - subtitle: subtitle text
      - author: author name
      - date: date string
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Gradient background
    _apply_bg_gradient(slide, theme["gradient_start"], theme["gradient_end"], angle=225)

    # Decorative top-right circle (semi-transparent look via lighter color)
    _add_circle(slide, 10.5, -1.5, 4, _lighten(theme["accent"], 0.2))
    # Decorative bottom-left circle
    _add_circle(slide, -1.5, 5.5, 3.5, _darken(theme["primary"], 0.5))

    # Accent line
    _add_decorative_bar(slide, 1.5, 3.2, 2.5, 0.06, theme["accent"])

    # Title
    title_text = slide_data.get("title", "Presentation Title")
    _add_text_box(slide, 1.5, 1.5, 10, 1.8, {
        "text": title_text, "font_size": 42, "bold": True,
        "color": theme["text_light"], "align": "LEFT",
        "vertical_anchor": "BOTTOM", "font_name": default_font,
    }, default_font)

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_text_box(slide, 1.5, 3.5, 9, 1.0, {
            "text": subtitle, "font_size": 20,
            "color": _lighten(theme["text_light"], 0.15), "align": "LEFT",
            "vertical_anchor": "TOP",
        }, default_font)

    # Author & Date
    author = slide_data.get("author", "")
    date = slide_data.get("date", "")
    meta_text = "  |  ".join(filter(None, [author, date]))
    if meta_text:
        _add_text_box(slide, 1.5, 5.8, 8, 0.5, {
            "text": meta_text, "font_size": 14,
            "color": _lighten(theme["text_light"], 0.3), "align": "LEFT",
        }, default_font)

    # Bottom accent bar
    _add_decorative_bar(slide, 0, 7.25, 13.333, 0.08, theme["accent"])

    # Slide number placeholder (hidden on title)
    return slide


def _build_section_slide(prs, slide_data, theme, default_font):
    """Build a section divider slide with large centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_gradient(slide, theme["gradient_start"], _darken(theme["gradient_end"], 0.8), angle=200)

    # Decorative large number or icon
    section_num = slide_data.get("section_number", "")
    if section_num:
        _add_text_box(slide, 1, 1.5, 4, 3, {
            "text": str(section_num).zfill(2), "font_size": 96, "bold": True,
            "color": _lighten(theme["accent"], 0.15), "align": "LEFT",
            "vertical_anchor": "MIDDLE",
        }, default_font)

    # Section title
    title = slide_data.get("title", "Section Title")
    _add_text_box(slide, 1.5, 3.0, 10, 2.0, {
        "text": title, "font_size": 38, "bold": True,
        "color": theme["text_light"], "align": "LEFT",
        "vertical_anchor": "MIDDLE",
    }, default_font)

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_decorative_bar(slide, 1.5, 5.0, 2, 0.05, theme["accent"])
        _add_text_box(slide, 1.5, 5.2, 9, 0.8, {
            "text": subtitle, "font_size": 18,
            "color": _lighten(theme["text_light"], 0.2), "align": "LEFT",
        }, default_font)

    _add_decorative_bar(slide, 0, 7.25, 13.333, 0.08, theme["accent"])
    return slide


def _build_content_slide(prs, slide_data, theme, default_font):
    """
    Build a standard content slide with title bar and body content.
    slide_data:
      - title: slide title
      - content: text content config (supports text, paragraphs, bullets)
      - elements: optional list of additional elements
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar background
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])

    # Title text
    title = slide_data.get("title", "")
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": title, "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT",
        "vertical_anchor": "MIDDLE",
    }, default_font)

    # Accent underline
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    # Body content
    content = slide_data.get("content")
    if content:
        content_cfg = content if isinstance(content, dict) else {"text": str(content)}
        content_cfg.setdefault("font_size", 18)
        content_cfg.setdefault("color", theme.get("text_dark", "#333333"))
        content_cfg.setdefault("line_spacing", 1.5)
        _add_text_box(slide, 0.8, 1.4, 11.5, 5.2, content_cfg, default_font)

    # Additional elements
    for elem in slide_data.get("elements", []):
        _add_element(slide, elem, theme, default_font)

    # Footer
    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_two_column_slide(prs, slide_data, theme, default_font):
    """
    Build a two-column layout slide.
    slide_data:
      - title: slide title
      - left_column: content config for left column
      - right_column: content config for right column
      - left_title / right_title: optional column titles
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    # Left column card
    _add_rounded_rect(slide, 0.5, 1.5, 5.9, 5.2, theme["card_bg"],
                      border_color=theme["card_border"], border_width=1.5, shadow=True)

    left_title = slide_data.get("left_title", "")
    if left_title:
        _add_decorative_bar(slide, 0.8, 1.7, 1.5, 0.05, theme["accent"])
        _add_text_box(slide, 0.8, 1.85, 5.3, 0.5, {
            "text": left_title, "font_size": 18, "bold": True,
            "color": theme["primary"], "align": "LEFT",
        }, default_font)
        left_content = slide_data.get("left_column", {})
        if isinstance(left_content, str):
            left_content = {"text": left_content}
        left_content.setdefault("font_size", 15)
        left_content.setdefault("color", theme.get("text_dark", "#333333"))
        _add_text_box(slide, 0.8, 2.5, 5.3, 3.8, left_content, default_font)
    else:
        left_content = slide_data.get("left_column", {})
        if isinstance(left_content, str):
            left_content = {"text": left_content}
        left_content.setdefault("font_size", 15)
        left_content.setdefault("color", theme.get("text_dark", "#333333"))
        _add_text_box(slide, 0.8, 1.8, 5.3, 4.5, left_content, default_font)

    # Right column card
    _add_rounded_rect(slide, 6.9, 1.5, 5.9, 5.2, theme["card_bg"],
                      border_color=theme["card_border"], border_width=1.5, shadow=True)

    right_title = slide_data.get("right_title", "")
    if right_title:
        _add_decorative_bar(slide, 7.2, 1.7, 1.5, 0.05, theme["accent"])
        _add_text_box(slide, 7.2, 1.85, 5.3, 0.5, {
            "text": right_title, "font_size": 18, "bold": True,
            "color": theme["primary"], "align": "LEFT",
        }, default_font)
        right_content = slide_data.get("right_column", {})
        if isinstance(right_content, str):
            right_content = {"text": right_content}
        right_content.setdefault("font_size", 15)
        right_content.setdefault("color", theme.get("text_dark", "#333333"))
        _add_text_box(slide, 7.2, 2.5, 5.3, 3.8, right_content, default_font)
    else:
        right_content = slide_data.get("right_column", {})
        if isinstance(right_content, str):
            right_content = {"text": right_content}
        right_content.setdefault("font_size", 15)
        right_content.setdefault("color", theme.get("text_dark", "#333333"))
        _add_text_box(slide, 7.2, 1.8, 5.3, 4.5, right_content, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_three_column_slide(prs, slide_data, theme, default_font):
    """Build a three-column card layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    columns = slide_data.get("columns", [])
    col_width = 3.7
    gap = 0.45
    start_x = 0.5

    for i, col in enumerate(columns[:3]):
        x = start_x + i * (col_width + gap)

        # Card background
        _add_rounded_rect(slide, x, 1.5, col_width, 5.2, theme["card_bg"],
                          border_color=theme["card_border"], border_width=1.5, shadow=True)

        # Card accent top bar
        _add_decorative_bar(slide, x + 0.3, 1.7, 1.2, 0.05, theme["accent"])

        # Column icon/number
        icon = col.get("icon", "")
        if icon:
            _add_text_box(slide, x + 0.3, 1.9, 1, 0.8, {
                "text": icon, "font_size": 28, "align": "LEFT",
                "color": theme["primary"], "bold": True,
            }, default_font)

        # Column title
        col_title = col.get("title", "")
        title_top = 2.0 if not icon else 2.7
        if col_title:
            _add_text_box(slide, x + 0.3, title_top, col_width - 0.6, 0.5, {
                "text": col_title, "font_size": 16, "bold": True,
                "color": theme["primary"], "align": "LEFT",
            }, default_font)

        # Column content
        col_content = col.get("content", {})
        if isinstance(col_content, str):
            col_content = {"text": col_content}
        col_content.setdefault("font_size", 13)
        col_content.setdefault("color", theme.get("text_dark", "#333333"))
        content_top = title_top + 0.6
        _add_text_box(slide, x + 0.3, content_top, col_width - 0.6, 5.2 - (content_top - 1.5) + 0.3,
                      col_content, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_cards_slide(prs, slide_data, theme, default_font):
    """
    Build a slide with N cards (2-6) arranged in a grid.
    slide_data:
      - title: slide title
      - cards: [{icon, title, content}, ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    cards = slide_data.get("cards", [])
    n = len(cards)
    if n == 0:
        _add_slide_footer(slide, prs, theme, default_font)
        return slide

    # Determine grid layout
    if n <= 3:
        cols_count, rows_count = n, 1
    elif n <= 4:
        cols_count, rows_count = 2, 2
    elif n <= 6:
        cols_count, rows_count = 3, 2
    else:
        cols_count, rows_count = 3, 2
        cards = cards[:6]
        n = 6

    total_w = 12.3
    total_h = 5.0 if rows_count == 1 else 4.8
    gap = 0.35
    card_w = (total_w - (cols_count - 1) * gap) / cols_count
    card_h = (total_h - (rows_count - 1) * gap) / rows_count
    start_x = 0.5
    start_y = 1.4

    for idx, card in enumerate(cards):
        row = idx // cols_count
        col = idx % cols_count
        cx = start_x + col * (card_w + gap)
        cy = start_y + row * (card_h + gap)

        # Card background
        _add_rounded_rect(slide, cx, cy, card_w, card_h, theme["card_bg"],
                          border_color=theme["card_border"], border_width=1.5, shadow=True)

        # Accent top bar on card
        _add_decorative_bar(slide, cx + 0.25, cy + 0.15, 0.8, 0.04, theme["accent"])

        # Icon
        icon = card.get("icon", "")
        current_y = cy + 0.3
        if icon:
            _add_text_box(slide, cx + 0.25, current_y, 1, 0.6, {
                "text": icon, "font_size": 22, "bold": True,
                "color": theme["primary"], "align": "LEFT",
            }, default_font)
            current_y += 0.55

        # Card title
        card_title = card.get("title", "")
        if card_title:
            _add_text_box(slide, cx + 0.25, current_y, card_w - 0.5, 0.4, {
                "text": card_title, "font_size": 14, "bold": True,
                "color": theme["primary"], "align": "LEFT",
            }, default_font)
            current_y += 0.45

        # Card content
        card_content = card.get("content", "")
        if card_content:
            if isinstance(card_content, str):
                card_content = {"text": card_content}
            card_content.setdefault("font_size", 12)
            card_content.setdefault("color", theme.get("text_dark", "#555555"))
            card_content.setdefault("line_spacing", 1.3)
            remaining_h = card_h - (current_y - cy) - 0.2
            _add_text_box(slide, cx + 0.25, current_y, card_w - 0.5, remaining_h,
                          card_content, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_chart_slide(prs, slide_data, theme, default_font):
    """
    Build a slide with a native PowerPoint chart.
    slide_data:
      - title: slide title
      - chart: chart configuration dict
      - description: optional text description below/beside chart
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    chart_config = slide_data.get("chart", {})
    description = slide_data.get("description", "")

    if description:
        # Chart on left, description on right
        _add_native_chart(slide, chart_config, left=0.5, top=1.4, width=8, height=5.2)
        _add_rounded_rect(slide, 8.8, 1.5, 4.2, 5.0, theme["card_bg"],
                          border_color=theme["card_border"], border_width=1)
        desc_cfg = description if isinstance(description, dict) else {"text": str(description)}
        desc_cfg.setdefault("font_size", 14)
        desc_cfg.setdefault("color", theme.get("text_dark", "#333333"))
        desc_cfg.setdefault("line_spacing", 1.4)
        _add_text_box(slide, 9.0, 1.7, 3.8, 4.6, desc_cfg, default_font)
    else:
        # Full-width chart
        _add_native_chart(slide, chart_config, left=0.8, top=1.4, width=11.5, height=5.3)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_stats_slide(prs, slide_data, theme, default_font):
    """
    Build a slide showcasing key statistics/numbers.
    slide_data:
      - title: slide title
      - stats: [{value: "98%", label: "Accuracy", icon: "..."}, ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    stats = slide_data.get("stats", [])
    n = len(stats)
    if n == 0:
        _add_slide_footer(slide, prs, theme, default_font)
        return slide

    total_w = 12.3
    gap = 0.4
    card_w = (total_w - (n - 1) * gap) / n if n <= 5 else (total_w - 4 * gap) / 5
    start_x = 0.5
    card_h = 3.5
    card_y = 2.2

    for i, stat in enumerate(stats[:5]):
        cx = start_x + i * (card_w + gap)

        # Card
        _add_rounded_rect(slide, cx, card_y, card_w, card_h, theme["card_bg"],
                          border_color=theme["card_border"], border_width=1.5, shadow=True)

        # Top accent
        _add_decorative_bar(slide, cx, card_y, card_w, 0.06, theme["accent"])

        # Icon
        icon = stat.get("icon", "")
        current_y = card_y + 0.3
        if icon:
            _add_text_box(slide, cx + 0.2, current_y, card_w - 0.4, 0.6, {
                "text": icon, "font_size": 24, "align": "CENTER",
                "color": theme["accent"],
            }, default_font)
            current_y += 0.6

        # Big number/value
        value = stat.get("value", "0")
        _add_text_box(slide, cx + 0.2, current_y, card_w - 0.4, 1.0, {
            "text": str(value), "font_size": 36, "bold": True,
            "color": theme["primary"], "align": "CENTER",
            "vertical_anchor": "MIDDLE",
        }, default_font)
        current_y += 1.0

        # Label
        label = stat.get("label", "")
        if label:
            _add_text_box(slide, cx + 0.2, current_y, card_w - 0.4, 0.6, {
                "text": label, "font_size": 14,
                "color": theme.get("text_dark", "#555555"), "align": "CENTER",
            }, default_font)

        # Description
        desc = stat.get("description", "")
        if desc:
            _add_text_box(slide, cx + 0.2, current_y + 0.55, card_w - 0.4, 0.8, {
                "text": desc, "font_size": 11,
                "color": "#888888", "align": "CENTER", "line_spacing": 1.2,
            }, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_timeline_slide(prs, slide_data, theme, default_font):
    """
    Build a horizontal timeline slide.
    slide_data:
      - title: slide title
      - steps: [{title, description, time_label}, ...]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    steps = slide_data.get("steps", [])
    n = len(steps)
    if n == 0:
        _add_slide_footer(slide, prs, theme, default_font)
        return slide

    # Timeline horizontal line
    line_y = 3.8
    _add_decorative_bar(slide, 0.8, line_y, 11.7, 0.04, theme["secondary"])

    step_width = 11.7 / n
    start_x = 0.8

    for i, step in enumerate(steps):
        cx = start_x + i * step_width + step_width / 2

        # Circle node on timeline
        node_size = 0.4
        _add_circle(slide, cx - node_size / 2, line_y - node_size / 2 + 0.02,
                    node_size, theme["accent"], border_color=theme["primary"], border_width=2)

        # Step number inside circle
        _add_text_box(slide, cx - node_size / 2, line_y - node_size / 2 + 0.02,
                      node_size, node_size, {
                          "text": str(i + 1), "font_size": 12, "bold": True,
                          "color": theme["text_light"], "align": "CENTER",
                          "vertical_anchor": "MIDDLE", "margin_left": 0, "margin_right": 0,
                          "margin_top": 0, "margin_bottom": 0,
                      }, default_font)

        # Alternate above/below
        if i % 2 == 0:
            # Content above the line
            time_label = step.get("time_label", "")
            if time_label:
                _add_text_box(slide, cx - step_width / 2 + 0.1, 1.5, step_width - 0.2, 0.4, {
                    "text": time_label, "font_size": 11, "bold": True,
                    "color": theme["accent"], "align": "CENTER",
                }, default_font)

            step_title = step.get("title", "")
            _add_text_box(slide, cx - step_width / 2 + 0.1, 1.9, step_width - 0.2, 0.5, {
                "text": step_title, "font_size": 14, "bold": True,
                "color": theme["primary"], "align": "CENTER",
            }, default_font)

            desc = step.get("description", "")
            if desc:
                _add_text_box(slide, cx - step_width / 2 + 0.1, 2.4, step_width - 0.2, 1.2, {
                    "text": desc, "font_size": 11,
                    "color": theme.get("text_dark", "#555555"), "align": "CENTER",
                    "line_spacing": 1.2,
                }, default_font)
        else:
            # Content below the line
            time_label = step.get("time_label", "")
            step_title = step.get("title", "")
            desc = step.get("description", "")

            _add_text_box(slide, cx - step_width / 2 + 0.1, 4.2, step_width - 0.2, 0.5, {
                "text": step_title, "font_size": 14, "bold": True,
                "color": theme["primary"], "align": "CENTER",
            }, default_font)

            if time_label:
                _add_text_box(slide, cx - step_width / 2 + 0.1, 4.7, step_width - 0.2, 0.4, {
                    "text": time_label, "font_size": 11, "bold": True,
                    "color": theme["accent"], "align": "CENTER",
                }, default_font)

            if desc:
                _add_text_box(slide, cx - step_width / 2 + 0.1, 5.1, step_width - 0.2, 1.2, {
                    "text": desc, "font_size": 11,
                    "color": theme.get("text_dark", "#555555"), "align": "CENTER",
                    "line_spacing": 1.2,
                }, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_table_slide(prs, slide_data, theme, default_font):
    """Build a slide with a professionally styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    table_config = slide_data.get("table", {})
    table_config.setdefault("left", 0.5)
    table_config.setdefault("top", 1.5)
    table_config.setdefault("width", 12.3)
    _add_styled_table(slide, table_config, theme, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_image_slide(prs, slide_data, theme, default_font, workspace):
    """
    Build a slide with an image and optional caption/description.
    slide_data:
      - title: slide title
      - image_path: path to image file
      - caption: image caption text
      - description: text description beside the image
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    img_path = slide_data.get("image_path", "")
    if img_path:
        full_path = img_path if os.path.isabs(img_path) else os.path.join(workspace, img_path)
        if os.path.exists(full_path):
            description = slide_data.get("description", "")
            if description:
                # Image on left, description on right
                slide.shapes.add_picture(full_path, _inches(0.8), _inches(1.5),
                                         _inches(7.5), _inches(5))
                desc_cfg = description if isinstance(description, dict) else {"text": str(description)}
                desc_cfg.setdefault("font_size", 14)
                desc_cfg.setdefault("color", theme.get("text_dark", "#333333"))
                _add_text_box(slide, 8.8, 1.5, 4, 5, desc_cfg, default_font)
            else:
                # Centered image
                slide.shapes.add_picture(full_path, _inches(1.5), _inches(1.5),
                                         _inches(10), _inches(5))

    caption = slide_data.get("caption", "")
    if caption:
        _add_text_box(slide, 0.8, 6.5, 11.5, 0.5, {
            "text": caption, "font_size": 12, "italic": True,
            "color": "#888888", "align": "CENTER",
        }, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_comparison_slide(prs, slide_data, theme, default_font):
    """
    Build a comparison slide (e.g., Before vs After, Pros vs Cons).
    slide_data:
      - title: slide title
      - left_title, right_title: column headers
      - left_items, right_items: lists of items
      - left_color, right_color: optional custom colors
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_solid(slide, theme.get("bg", "#FFFFFF"))

    # Title bar
    _add_rounded_rect(slide, 0, 0, 13.333, 1.1, theme["primary"])
    _add_text_box(slide, 0.8, 0.15, 11, 0.8, {
        "text": slide_data.get("title", ""), "font_size": 26, "bold": True,
        "color": theme["text_light"], "align": "LEFT", "vertical_anchor": "MIDDLE",
    }, default_font)
    _add_decorative_bar(slide, 0.8, 1.05, 3, 0.04, theme["accent"])

    left_color = slide_data.get("left_color", theme["primary"])
    right_color = slide_data.get("right_color", theme["accent"])

    # Left card
    _add_rounded_rect(slide, 0.5, 1.5, 5.9, 5.2, theme["card_bg"],
                      border_color=left_color, border_width=2, shadow=True)
    _add_decorative_bar(slide, 0.5, 1.5, 5.9, 0.06, left_color)

    left_title = slide_data.get("left_title", "Option A")
    _add_text_box(slide, 0.8, 1.7, 5.3, 0.6, {
        "text": left_title, "font_size": 20, "bold": True,
        "color": left_color, "align": "CENTER",
    }, default_font)

    left_items = slide_data.get("left_items", [])
    if left_items:
        _add_text_box(slide, 0.8, 2.5, 5.3, 3.8, {
            "bullets": left_items, "font_size": 14,
            "color": theme.get("text_dark", "#333333"), "line_spacing": 1.5,
        }, default_font)

    # VS divider
    _add_circle(slide, 6.1, 3.5, 0.8, theme["secondary"])
    _add_text_box(slide, 6.1, 3.5, 0.8, 0.8, {
        "text": "VS", "font_size": 14, "bold": True,
        "color": theme["text_light"], "align": "CENTER",
        "vertical_anchor": "MIDDLE", "margin_left": 0, "margin_right": 0,
        "margin_top": 0, "margin_bottom": 0,
    }, default_font)

    # Right card
    _add_rounded_rect(slide, 6.9, 1.5, 5.9, 5.2, theme["card_bg"],
                      border_color=right_color, border_width=2, shadow=True)
    _add_decorative_bar(slide, 6.9, 1.5, 5.9, 0.06, right_color)

    right_title = slide_data.get("right_title", "Option B")
    _add_text_box(slide, 7.2, 1.7, 5.3, 0.6, {
        "text": right_title, "font_size": 20, "bold": True,
        "color": right_color, "align": "CENTER",
    }, default_font)

    right_items = slide_data.get("right_items", [])
    if right_items:
        _add_text_box(slide, 7.2, 2.5, 5.3, 3.8, {
            "bullets": right_items, "font_size": 14,
            "color": theme.get("text_dark", "#333333"), "line_spacing": 1.5,
        }, default_font)

    _add_slide_footer(slide, prs, theme, default_font)
    return slide


def _build_quote_slide(prs, slide_data, theme, default_font):
    """Build a quote/highlight slide with large centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_gradient(slide, theme["gradient_start"], theme["gradient_end"], angle=225)

    # Large quote mark decoration
    _add_text_box(slide, 1, 1.2, 3, 2, {
        "text": "\u201C", "font_size": 120, "bold": True,
        "color": _lighten(theme["accent"], 0.2), "align": "LEFT",
        "vertical_anchor": "TOP",
    }, default_font)

    # Quote text
    quote = slide_data.get("quote", slide_data.get("text", ""))
    _add_text_box(slide, 2, 2.5, 9.5, 3, {
        "text": quote, "font_size": 28, "italic": True,
        "color": theme["text_light"], "align": "CENTER",
        "vertical_anchor": "MIDDLE", "line_spacing": 1.5,
    }, default_font)

    # Author
    author = slide_data.get("author", "")
    if author:
        _add_decorative_bar(slide, 5.5, 5.5, 2, 0.04, theme["accent"])
        _add_text_box(slide, 2, 5.7, 9.5, 0.6, {
            "text": f"\u2014 {author}", "font_size": 18,
            "color": _lighten(theme["text_light"], 0.2), "align": "CENTER",
        }, default_font)

    _add_decorative_bar(slide, 0, 7.25, 13.333, 0.08, theme["accent"])
    return slide


def _build_ending_slide(prs, slide_data, theme, default_font):
    """Build a thank-you / ending slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg_gradient(slide, theme["gradient_start"], theme["gradient_end"], angle=225)

    # Decorative circles
    _add_circle(slide, 10, -1, 3.5, _lighten(theme["accent"], 0.15))
    _add_circle(slide, -1, 5, 3, _darken(theme["primary"], 0.5))

    # Main text
    main_text = slide_data.get("title", slide_data.get("text", "Thank You"))
    _add_text_box(slide, 1, 2.0, 11.3, 2.5, {
        "text": main_text, "font_size": 48, "bold": True,
        "color": theme["text_light"], "align": "CENTER",
        "vertical_anchor": "MIDDLE",
    }, default_font)

    # Subtitle
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        _add_text_box(slide, 2, 4.5, 9.3, 1.0, {
            "text": subtitle, "font_size": 20,
            "color": _lighten(theme["text_light"], 0.2), "align": "CENTER",
        }, default_font)

    # Contact info
    contact = slide_data.get("contact", "")
    if contact:
        _add_text_box(slide, 2, 5.5, 9.3, 0.8, {
            "text": contact, "font_size": 16,
            "color": _lighten(theme["text_light"], 0.3), "align": "CENTER",
        }, default_font)

    _add_decorative_bar(slide, 0, 7.25, 13.333, 0.08, theme["accent"])
    return slide


# ============================================================================
#  ELEMENT DISPATCHER — For adding arbitrary elements to any slide
# ============================================================================

def _add_element(slide, elem, theme, default_font):
    """Add an arbitrary element to a slide (used by content slides)."""
    etype = elem.get("type", "text")

    if etype == "text":
        return _add_text_box(slide, elem.get("left", 1), elem.get("top", 2),
                             elem.get("width", 10), elem.get("height", 1),
                             elem, default_font)

    elif etype == "shape":
        shape_type = elem.get("shape_type", "ROUNDED_RECTANGLE")
        mso = getattr(MSO_SHAPE, shape_type.upper(), MSO_SHAPE.ROUNDED_RECTANGLE)
        shape = slide.shapes.add_shape(
            mso, _inches(elem.get("left", 1)), _inches(elem.get("top", 2)),
            _inches(elem.get("width", 3)), _inches(elem.get("height", 1))
        )
        color = elem.get("color", theme["accent"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex(color)
        if elem.get("border_color"):
            shape.line.color.rgb = _hex(elem["border_color"])
            shape.line.width = Pt(elem.get("border_width", 1))
        else:
            shape.line.fill.background()
        if "text" in elem:
            _add_text_to_frame(shape.text_frame, elem, default_font)
        if "rotation" in elem:
            shape.rotation = elem["rotation"]
        return shape

    elif etype == "image":
        workspace = get_workspace_path() or "."
        img_path = elem.get("path", "")
        full_path = img_path if os.path.isabs(img_path) else os.path.join(workspace, img_path)
        if os.path.exists(full_path):
            return slide.shapes.add_picture(
                full_path, _inches(elem.get("left", 1)), _inches(elem.get("top", 2)),
                _inches(elem.get("width", 5)), _inches(elem.get("height", 3))
            )

    elif etype == "chart":
        return _add_native_chart(
            slide, elem,
            left=elem.get("left", 1), top=elem.get("top", 2),
            width=elem.get("width", 8), height=elem.get("height", 4.5)
        )

    elif etype == "table":
        return _add_styled_table(slide, elem, theme, default_font)

    return None


# ============================================================================
#  FOOTER — Consistent slide footer with page number
# ============================================================================

def _add_slide_footer(slide, prs, theme, default_font):
    """Add a consistent footer bar and slide number."""
    # Bottom accent bar
    _add_decorative_bar(slide, 0, 7.2, 13.333, 0.04, theme["accent"])

    # Slide number
    slide_num = len(prs.slides)
    _add_text_box(slide, 12.2, 7.0, 0.8, 0.3, {
        "text": str(slide_num), "font_size": 10,
        "color": theme.get("secondary", "#888888"), "align": "RIGHT",
        "vertical_anchor": "MIDDLE",
    }, default_font)


# ============================================================================
#  MAIN ENTRY POINT
# ============================================================================

SLIDE_BUILDERS = {
    "title": _build_title_slide,
    "cover": _build_title_slide,
    "section": _build_section_slide,
    "divider": _build_section_slide,
    "content": _build_content_slide,
    "text": _build_content_slide,
    "two_column": _build_two_column_slide,
    "two_columns": _build_two_column_slide,
    "three_column": _build_three_column_slide,
    "three_columns": _build_three_column_slide,
    "cards": _build_cards_slide,
    "chart": _build_chart_slide,
    "stats": _build_stats_slide,
    "statistics": _build_stats_slide,
    "timeline": _build_timeline_slide,
    "table": _build_table_slide,
    "image": _build_image_slide,
    "comparison": _build_comparison_slide,
    "quote": _build_quote_slide,
    "ending": _build_ending_slide,
    "thank_you": _build_ending_slide,
}


def generate_pptx_presentation(file_path, slides_content, presentation_settings=None):
    """
    Ultra PPT Engine v2.0 — Generate professional presentations with rich visual design.

    Args:
        file_path: Output .pptx file path (relative to workspace)
        slides_content: List of slide definitions, each with a "slide_type" field
        presentation_settings: Global settings (theme, default_font, title, author, etc.)

    Returns:
        dict with status and file_path
    """
    workspace = get_workspace_path()
    if not workspace:
        return {"error": "Workspace not configured. Please set up workspace first."}

    full_path = os.path.join(workspace, file_path)
    os.makedirs(os.path.dirname(full_path) if os.path.dirname(full_path) else workspace, exist_ok=True)

    settings = presentation_settings or {}
    theme_name = settings.get("theme", "midnight")
    theme = _get_theme(theme_name)

    # Allow custom theme color overrides
    custom_colors = settings.get("custom_colors", {})
    if custom_colors:
        theme = {**theme, **custom_colors}

    default_font = settings.get("default_font", "微软雅黑")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Metadata
    prs.core_properties.author = settings.get("author", "Ultra PPT Engine v2.0")
    prs.core_properties.title = settings.get("title", "Professional Presentation")

    # Build each slide
    for slide_data in slides_content:
        slide_type = slide_data.get("slide_type", "content").lower().strip()
        builder = SLIDE_BUILDERS.get(slide_type)

        if builder:
            if slide_type in ("image",):
                builder(prs, slide_data, theme, default_font, workspace)
            else:
                builder(prs, slide_data, theme, default_font)
        else:
            # Fallback to content slide
            _build_content_slide(prs, slide_data, theme, default_font)

    prs.save(full_path)
    return {
        "success": True,
        "file_path": full_path,
        "slides_count": len(prs.slides),
        "theme": theme_name,
        "message": f"Presentation saved with {len(prs.slides)} slides using '{theme_name}' theme."
    }
